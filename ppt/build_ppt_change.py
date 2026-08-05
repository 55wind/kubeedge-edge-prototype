"""ETRI 2차 킥오프 매핑 자료 기준 '기존(1차년도) -> 현재(2차년도)' 화면 구성 덱.

수행항목 1~7 각각을 한 페이지로 구성한다. 페이지마다
(1) 기존 -> 현재 변화 요약 밴드, (2) 그 변화가 실제로 담긴 실행 화면 캡처만 싣는다.
표지/아키텍처/향후계획 등 화면이 없는 페이지는 넣지 않는다 (요구사항).

캡처 원본은 모두 docs/screens/*.png 이며, 각 재생성 스크립트는 다음과 같다:
  ppt/capture_platform_screens.py - 실행 중인 Manager를 Swagger UI에서 직접 호출한 화면
  ppt/capture_k8s_screens.py      - 실제 Kubernetes 클러스터 상태·Dashboard·파드 로그
  ppt/capture_cli_screens.py      - 실제 CLI 실행 출력(demo/gateway/fleet)
  ppt/capture_screens.py          - 문서·매니페스트 대조 화면 2종
성능 차트: docs/perf/*.png (bench/report.py 로 재생성).

실행:
    python ppt/build_ppt_change.py            # 생성 + 검증
    python ppt/build_ppt_change.py --verify   # 검증만
"""
from __future__ import annotations

import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from build_ppt import (  # noqa: E402  (동일 디렉터리 헬퍼 재사용)
    GRAY_ACCENT,
    GRAY_BODY,
    NAVY_DARK,
    NAVY_MID,
    P95_CHART,
    SCREEN_AUDIT,
    SCREEN_DEMO,
    SCREEN_SWAGGER,
    SCREENS_DIR,
    SLIDE_H,
    SLIDE_W,
    WHITE,
    _set_run_font,
    add_footer,
    add_rect,
    add_screen_picture,
    add_slide,
    add_title_bar,
)

REPO_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_PATH = REPO_ROOT / "ppt" / "ETRI_2차년도_기존-현재_화면구성.pptx"
SCREEN_COMM = SCREENS_DIR / "screen_comm_decision.png"
SCREEN_K8S_YAML = SCREENS_DIR / "screen_k8s_yaml.png"
SCREEN_GATEWAY = SCREENS_DIR / "screen_gateway.png"
# 실제 kind 클러스터 / 실행 중인 Manager의 Swagger UI 캡처
SCREEN_K8S_PODS = SCREENS_DIR / "k8s_pods.png"
SCREEN_K8S_DASHBOARD = SCREENS_DIR / "k8s_dashboard.png"
SCREEN_K8S_AGENT_LOG = SCREENS_DIR / "k8s_agent_log.png"
SCREEN_REGISTER = SCREENS_DIR / "platform_register.png"
SCREEN_RBAC_DENIED = SCREENS_DIR / "platform_rbac_denied.png"

# (번호, 제목, 기존 불릿, 현재 불릿, [(이미지, 캡션), ...])
CHANGE_SLIDES = [
    (
        1, "쿠버네티스(KubeEdge) 기반 구동·성능",
        [
            "KubeEdge 클러스터 구축 스크립트까지만 - 보안 기능의 K8s 구동·성능 검증 없음",
            "앱 이미지는 저장소 밖에서 수동 빌드 (소스 미포함)",
        ],
        [
            "보안 코어 모듈이 실제 K8s 클러스터에서 구동 - manager/agent×2/gateway 4개 파드 Running",
            "가상 디바이스 실측(N=10~200) + 1,000기 M/M/c 외삽: ρ=0.70 안정, 대기 p95 372.6ms, 권장 replica 1",
        ],
        [
            (SCREEN_K8S_PODS, "실제 클러스터 상태 - edge 라벨 노드에 agent/gateway 배치, manager NodePort 30443"),
            (P95_CHART, "실측 auth p95 지연 vs N (bench/run_bench.py 산출)"),
        ],
    ),
    (
        2, "미들웨어 통신 방식 결정",
        [
            "채널(Channel) 기반 vs 별도 프로토콜 - 비교·선정 없음",
            "RabbitMQ 단일 큐 + 평문 자격증명 하드코딩",
        ],
        [
            "하이브리드 확정: 제어=KubeEdge 채널 / 데이터·인증=별도 mTLS HTTPS(REST)",
            "4축 비교표·판단 기준 수립 + 매니저 공존·샤딩·우선순위 규칙 문서화 (docs/01·02), AMQPS는 옵션 분리",
        ],
        [(SCREEN_COMM, "docs/01 통신 방식 비교 기준표와 하이브리드 선정 결론 (구현으로 실증됨)")],
    ),
    (
        3, "대규모(1,000기) 디바이스 인증 (AAA) ★핵심",
        [
            "단순 JWT 발급 + hostPath 인증서 수동 배포",
            "위조·폐기 대응 없음, 대규모(1,000기) 인증 검증 없음",
        ],
        [
            "등록(CSR SAN 신원대조) → X.509 발급 → JWT/RBAC → JWS 검증 → 감사 → 폐기 전 주기 구현",
            "위조 JWS 거부·폐기 즉시 차단이 감사로그에 실기록 + 실측 기반 1,000기 외삽 검증",
        ],
        [
            (SCREEN_REGISTER, "Authentication - 등록 요청에 X.509 인증서가 실제 발급된 응답"),
            (SCREEN_RBAC_DENIED, "Authorization - operator 역할의 admin 전용 API 호출이 403으로 거부"),
            (SCREEN_AUDIT, "Accounting - 인증 성공/위조 거부/폐기가 감사로그에 실기록"),
        ],
    ),
    (
        4, "KubeEdge 연동 구조 적용",
        [
            "CloudCore-EdgeCore 설치까지만 - 보안 모듈 적용 구조 없음",
            "IP·자격증명 하드코딩, nodeName으로 노드 직접 지정",
        ],
        [
            "4계층(Device-Agent-Manager-Backend) 보안 모듈을 KubeEdge 배포 구조에 적용",
            "edge 라벨 스케줄 + Secret 자동 생성·참조 + __CLOUD_IP__ 치환 (deploy/demo-setup-v2.sh)",
        ],
        [
            (SCREEN_K8S_DASHBOARD, "Kubernetes Dashboard - edge-auth 네임스페이스에 4개 Deployment 구동"),
            (SCREEN_K8S_AGENT_LOG, "agent 파드 로그 - 실제 등록→토큰→텔레메트리 사이클 성공"),
            (SCREEN_K8S_YAML, "매니페스트 1차년도(좌) vs 2차년도(우) - 하드코딩 제거, Secret·edge 라벨"),
        ],
    ),
    (
        5, "네트워크 구성·주소 체계 (공인 IP)",
        [
            "공인 IP 노출 위험 검토 없음",
            "디바이스가 Manager에 개별 직접 연결 (공인 접점 = 디바이스 수만큼)",
        ],
        [
            "게이트웨이+사설IP 구조 권고 확정 - 보안성·관리·확장성·1,000기 적용성 4축 비교 (docs/03)",
            "EdgeGateway 구현: 사설망 하위 디바이스를 단일 인증 업링크(JWS 배치)로 집선, 공인 접점 1개로 축소",
        ],
        [(SCREEN_GATEWAY, "EdgeGateway 실제 실행 - 하위 3기 배치 업링크, Manager에는 게이트웨이 신원만 노출")],
    ),
    (
        6, "보안 모듈-프레임워크 연동 역할 정립",
        [
            "보안 모듈과 프레임워크 간 통합 역할·인터페이스 불명확",
            "타 기관 연동 구조·범위 미정",
        ],
        [
            "상명대 = 보안 코어 모듈 제공자 확정 + R&R 문서화 (docs/04)",
            "API 계약 9종·인증서 프로파일·JWT 클레임 규격 명세, Swagger(OpenAPI)로 기계가독 계약 제공",
        ],
        [(SCREEN_SWAGGER, "Swagger UI(/docs) - Manager 전체 API 계약, docs/04 §2 인터페이스 명세와 1:1 대응")],
    ),
    (
        7, "시연 구성",
        [
            "메시지 왕복 데모(demo-test.sh)뿐 - 보안 적용 전·후 비교 없음",
        ],
        [
            "보안·인증 적용 전·후 비교 시연 구현: 미인증 주입 200 수용(BEFORE) → 401 거부(AFTER)",
            "정식 등록·인증 디바이스만 전송 성공, 감사로그 tail 동시 출력 (demo/run_demo.py, 영상 촬영 포인트 포함)",
        ],
        [(SCREEN_DEMO, "demo/run_demo.py --fast 실제 실행 - 전·후 비교와 감사로그가 한 화면에 기록")],
    ),
]


def _comparison_band(slide, before_lines, after_lines):
    """기존 -> 현재 변화 요약 밴드 (좌: 기존/회색, 우: 현재/남색, 가운데 화살표)."""
    top = Inches(1.08)
    h = Inches(1.52)

    before_box = add_rect(slide, Inches(0.5), top, Inches(5.55), h, fill=GRAY_ACCENT)
    tf = before_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.07)
    head = tf.paragraphs[0].add_run()
    head.text = "기존 (1차년도)"
    _set_run_font(head, size=12.5, bold=True, color=NAVY_DARK)
    for line in before_lines:
        p = tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = f"• {line}"
        _set_run_font(r, size=11, color=GRAY_BODY)

    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(6.13), top + Inches(0.5), Inches(0.62), Inches(0.52)
    )
    arrow.shadow.inherit = False
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = NAVY_DARK
    arrow.line.fill.background()

    after_box = add_rect(slide, Inches(6.85), top, Inches(5.98), h, fill=NAVY_MID)
    tf2 = after_box.text_frame
    tf2.word_wrap = True
    tf2.vertical_anchor = MSO_ANCHOR.TOP
    tf2.margin_left = Inches(0.12)
    tf2.margin_right = Inches(0.1)
    tf2.margin_top = Inches(0.07)
    head2 = tf2.paragraphs[0].add_run()
    head2.text = "현재 (2차년도)"
    _set_run_font(head2, size=12.5, bold=True, color=WHITE)
    for line in after_lines:
        p = tf2.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = f"• {line}"
        _set_run_font(r, size=11, color=WHITE)


CAPTION_H = Inches(0.62)


def _scaled_size(img_path, max_w, max_h):
    from PIL import Image as PILImage

    with PILImage.open(img_path) as im:
        w, h = im.size
    scale = min(max_w / w, max_h / h)
    return int(w * scale), int(h * scale)


def place_screen(slide, img_path, left, top, max_w, max_h, caption, caption_y):
    """Top-align the image in its column and put the caption on a shared baseline.

    Screens have different aspect ratios, so centring each one vertically would
    scatter the captions. Images hang from a common top edge and every caption in
    the row sits at ``caption_y`` (just below the tallest image), which keeps the
    row aligned without leaving dead space under the short ones.
    """
    pic_w, pic_h = _scaled_size(img_path, max_w, max_h)
    slide.shapes.add_picture(str(img_path), left + int((max_w - pic_w) / 2), top,
                             width=pic_w, height=pic_h)
    box = slide.shapes.add_textbox(left, caption_y, max_w, CAPTION_H)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = caption
    _set_run_font(run, size=10.5, color=GRAY_BODY)


def build_change_slide(prs, no, title, before, after, screens, page_no):
    slide = add_slide(prs)
    add_title_bar(slide, f"{no}. {title}",
                  "ETRI 2차 킥오프 매핑 자료 기준 - 기존(1차년도) → 현재(2차년도), "
                  "캡처는 실제 플랫폼 화면(Swagger UI·Kubernetes Dashboard·실제 CLI 출력)")
    _comparison_band(slide, before, after)

    img_top = Inches(2.78)
    img_h = Inches(4.32)
    n = len(screens)
    if n == 1:
        img_path, caption = screens[0]
        add_screen_picture(slide, img_path, Inches(0.5), img_top, Inches(12.33), img_h,
                           caption=caption)
    else:
        gap = Inches(0.2)
        col_w = int((Inches(12.33) - gap * (n - 1)) / n)
        avail_h = img_h - CAPTION_H
        tallest = max(_scaled_size(p, col_w, avail_h)[1] for p, _ in screens)
        caption_y = img_top + tallest + Inches(0.08)
        for i, (img_path, caption) in enumerate(screens):
            x = Inches(0.5) + i * (col_w + gap)
            place_screen(slide, img_path, x, img_top, col_w, avail_h, caption, caption_y)

    add_footer(slide, page_no)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for no, title, before, after, screens in CHANGE_SLIDES:
        build_change_slide(prs, no, title, before, after, screens, page_no=no)
    return prs


def verify(path: Path) -> None:
    with zipfile.ZipFile(path) as zf:
        assert zf.testzip() is None, "zip 무결성 오류"

    prs = Presentation(str(path))
    slides = list(prs.slides)
    assert len(slides) == 7, f"슬라이드 수 불일치: {len(slides)} (기대: 7)"

    for si, slide in enumerate(slides, start=1):
        pics = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
        assert pics, f"슬라이드 {si}에 화면 캡처가 없음 (요구: 모든 페이지에 화면 포함)"
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            assert shape.left >= -Emu(1000) and shape.top >= -Emu(1000)
            assert shape.left + (shape.width or 0) <= prs.slide_width + Emu(1000)
            assert shape.top + (shape.height or 0) <= prs.slide_height + Emu(1000)

    print(f"검증 통과: {len(slides)}개 슬라이드, 전 페이지 화면 캡처 포함, zip 무결성 OK, 경계 내 배치.")


def main():
    args = sys.argv[1:]
    if "--verify" not in args:
        prs = build_presentation()
        prs.save(str(OUTPUT_PATH))
        print(f"생성 완료: {OUTPUT_PATH}")
    if "--no-verify" not in args:
        verify(OUTPUT_PATH)


if __name__ == "__main__":
    main()
