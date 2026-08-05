"""ETRI AI_EDGE 2차년도 - 상명대 보안 코어 모듈 수행항목 매핑 PPT 생성 스크립트.

재실행 가능(deterministic)하게 `ppt/ETRI_2차년도_수행항목_매핑.pptx`를 생성한다.
텍스트는 아래 저장소 문서에 근거하며, 새로운 수치/사실을 창작하지 않는다:

- docs/REQUIREMENTS_MAPPING.md   (7개 수행항목 x 산출물 매핑, 근거 원본)
- docs/plans/2026-08-03-y2-security-core.md  (요구사항 원문 요약, 아키텍처)
- docs/01~06-*.md                 (수행항목별 결정/결론)
- docs/perf/PERFORMANCE_REPORT.md (성능 실측치 - 여기 수치만 인용)
- demo/DEMO_SCENARIO.md            (시연 전/후 비교 절차)
- README.md                        (1차->2차 Before/After 표)
- docs/screens/*.png               (실제 실행 화면 캡처 - ppt/capture_screens.py 로 재생성)

실행:
    python ppt/build_ppt.py            # pptx (재)생성
    python ppt/build_ppt.py --verify   # 생성 없이 기존 pptx 검증만 수행
    python ppt/build_ppt.py --no-verify  # 생성만 하고 검증 생략(디버그용)
"""

from __future__ import annotations

import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# 경로/상수
# ---------------------------------------------------------------------------

REPO_ROOT = Path(__file__).resolve().parent.parent
PPT_DIR = REPO_ROOT / "ppt"
OUTPUT_PATH = PPT_DIR / "ETRI_2차년도_수행항목_매핑.pptx"
PERF_DIR = REPO_ROOT / "docs" / "perf"
P95_CHART = PERF_DIR / "p95_latency_vs_n.png"
THROUGHPUT_CHART = PERF_DIR / "throughput_vs_n.png"

# 실행 화면 캡처 (ppt/capture_screens.py 로 재생성 - 실제 프로그램 출력)
SCREENS_DIR = REPO_ROOT / "docs" / "screens"
SCREEN_DEMO = SCREENS_DIR / "screen_demo.png"
# 실제 Manager를 Swagger UI에서 직접 실행해 얻은 화면 (ppt/capture_platform_screens.py)
SCREEN_DEVICES = SCREENS_DIR / "platform_devices.png"
SCREEN_AUDIT = SCREENS_DIR / "platform_audit.png"
SCREEN_SWAGGER = SCREENS_DIR / "platform_swagger_overview.png"
SCREEN_FLEET = SCREENS_DIR / "screen_fleet.png"
# 실제 kind 기반 Kubernetes 클러스터 캡처 (ppt/capture_k8s_screens.py)
SCREEN_K8S_DASHBOARD = SCREENS_DIR / "k8s_dashboard.png"
SCREEN_K8S_PODS = SCREENS_DIR / "k8s_pods.png"

FONT_NAME = "맑은 고딕"

# 팔레트 - 남색/회색 계열 단색(무지개 금지)
NAVY_DARK = RGBColor(0x1F, 0x38, 0x64)   # 제목/헤더
NAVY_MID = RGBColor(0x2E, 0x54, 0x90)    # 보조 헤더/강조 박스
NAVY_LIGHT = RGBColor(0x8E, 0xA9, 0xDB)  # 옅은 강조
GRAY_BODY = RGBColor(0x40, 0x40, 0x40)   # 본문 텍스트
GRAY_ACCENT = RGBColor(0xD9, 0xD9, 0xD9)  # 배경 강조
GRAY_LINE = RGBColor(0xA6, 0xA6, 0xA6)   # 선/테두리
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OK_NAVY = RGBColor(0x1F, 0x38, 0x64)     # "상태" 배지도 남색 계열로 통일

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

DATE_STR = "2026-08-03"
PRESENTER = "상명대학교"


# ---------------------------------------------------------------------------
# 헬퍼
# ---------------------------------------------------------------------------

def _set_run_font(run, *, size=14, bold=False, color=GRAY_BODY, italic=False):
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    # 동아시아 폰트 지정도 명시적으로 채워야 일부 뷰어에서 폰트가 정확히 적용된다.
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT_NAME)


def add_rect(slide, left, top, width, height, *, fill=None, line=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_width or Pt(1)
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    lines,
    *,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    word_wrap=True,
):
    """lines: [(text, size, bold, color, bullet_indent), ...] 또는 단순 문자열 리스트."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, size, bold, color = item, 14, False, GRAY_BODY
        else:
            text, size, bold, color = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        _set_run_font(run, size=size, bold=bold, color=color)
    return box


def add_section_block(slide, left, top, width, height, header_text, bullets, *, body_size=12):
    """헤더 바(남색) + 본문 불릿(회색) 블록. 요구/해결방안/산출물 섹션에 사용."""
    header_h = Inches(0.36)
    add_rect(slide, left, top, width, header_h, fill=NAVY_DARK)
    hb = slide.shapes.add_textbox(left, top, width, header_h)
    htf = hb.text_frame
    htf.margin_left = Inches(0.1)
    htf.margin_top = 0
    htf.margin_bottom = 0
    htf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = htf.paragraphs[0]
    run = p.add_run()
    run.text = header_text
    _set_run_font(run, size=14, bold=True, color=WHITE)

    body_top = top + header_h + Inches(0.04)
    body_h = height - header_h - Inches(0.04)
    add_rect(slide, left, body_top, width, body_h, fill=None, line=GRAY_LINE, line_width=Pt(0.75))
    bb = slide.shapes.add_textbox(left + Inches(0.12), body_top + Inches(0.06), width - Inches(0.24), body_h - Inches(0.12))
    btf = bb.text_frame
    btf.word_wrap = True
    btf.margin_left = 0
    btf.margin_right = 0
    btf.margin_top = 0
    btf.margin_bottom = 0
    for i, text in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = f"• {text}"
        _set_run_font(run, size=body_size, bold=False, color=GRAY_BODY)
    return bb


def add_slide(prs, layout_index=6):
    layout = prs.slide_layouts[layout_index]  # 6 = Blank
    return prs.slides.add_slide(layout)


def add_title_bar(slide, title_text, subtitle_text=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.95), fill=NAVY_DARK)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.08), SLIDE_W - Inches(1.0), Inches(0.8))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    _set_run_font(run, size=26, bold=True, color=WHITE)
    if subtitle_text:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle_text
        _set_run_font(run2, size=13, bold=False, color=NAVY_LIGHT)


def add_footer(slide, page_no):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.16), Inches(8.0), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"ETRI AI_EDGE 2차년도 - {PRESENTER} 보안 코어 모듈  |  {DATE_STR}"
    _set_run_font(run, size=9, bold=False, color=GRAY_LINE)

    pb = slide.shapes.add_textbox(SLIDE_W - Inches(1.0), Inches(7.16), Inches(0.6), Inches(0.3))
    ptf = pb.text_frame
    ptf.margin_left = 0
    ptf.margin_top = 0
    pp = ptf.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    prun = pp.add_run()
    prun.text = str(page_no)
    _set_run_font(prun, size=9, bold=False, color=GRAY_LINE)


# ---------------------------------------------------------------------------
# 슬라이드 1: 표지
# ---------------------------------------------------------------------------

def build_cover(prs):
    slide = add_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY_DARK)
    add_rect(slide, 0, Inches(4.55), SLIDE_W, Inches(0.06), fill=NAVY_LIGHT)

    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "ETRI AI_EDGE 2차년도"
    _set_run_font(run, size=22, bold=False, color=NAVY_LIGHT)
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "수행항목-해결방안 매핑 결과보고"
    _set_run_font(run2, size=36, bold=True, color=WHITE)
    p3 = tf.add_paragraph()
    run3 = p3.add_run()
    run3.text = "보안 코어 모듈(EAM: Edge Auth Manager) 구현 산출물 매핑"
    _set_run_font(run3, size=16, bold=False, color=GRAY_ACCENT)

    tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.85), Inches(11.3), Inches(1.6))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p4 = tf2.paragraphs[0]
    run4 = p4.add_run()
    run4.text = PRESENTER
    _set_run_font(run4, size=18, bold=True, color=WHITE)
    p5 = tf2.add_paragraph()
    run5 = p5.add_run()
    run5.text = DATE_STR
    _set_run_font(run5, size=14, bold=False, color=NAVY_LIGHT)
    p6 = tf2.add_paragraph()
    run6 = p6.add_run()
    run6.text = "저장소: logperch - 소스 포함 완전체 산출물"
    _set_run_font(run6, size=11, bold=False, color=GRAY_ACCENT)


# ---------------------------------------------------------------------------
# 슬라이드 2: 2차년도 수행 범위 재정리
# ---------------------------------------------------------------------------

def build_scope(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "2차년도 수행 범위 재정리", "핵심 보안 기능 중심으로 재집중 (ETRI 2026-06-17 킥오프 회의록 기준)")

    add_section_block(
        slide,
        Inches(0.5), Inches(1.15), Inches(6.0), Inches(3.0),
        "핵심 보안 4대 기능",
        [
            "디바이스 등록(Registration) - bootstrap token + CSR -> X.509 발급",
            "통신 인증(Authentication) - mTLS(클라이언트 인증서) + Bearer JWT(RS256)",
            "AAA - Authentication / Authorization(RBAC) / Accounting(감사로그)",
            "데이터 전송 보안 - TLS 채널 + JWS(RS256) 페이로드 서명",
        ],
        body_size=14,
    )

    add_section_block(
        slide,
        Inches(6.7), Inches(1.15), Inches(6.13), Inches(3.0),
        "★핵심 수행항목",
        [
            "수행항목 3: 대규모(1,000기) 디바이스 인증 AAA",
            "담당: 박건우 · 홍승우",
            "1,000기 동시 인증 스톰 SLA 실측 기반 외삽 검증(슬라이드 12)",
        ],
        body_size=14,
    )

    add_section_block(
        slide,
        Inches(0.5), Inches(4.35), Inches(12.33), Inches(2.35),
        "범위 조정 메모",
        [
            "장애 인지·자동 복구는 2차년도 필수 범위에서 제외하고, 시연용 최소 수준"
            "(예: 정식 인증 실패 401 거부 + 감사기록)으로만 다룬다.",
            "프레임워크(KubeEdge/ETRI 플랫폼) 전체 통합은 3차년도 이관 또는 타 기관 연계 "
            "항목으로 조정하고, 본 저장소는 보안 코어 모듈의 REST API 계약(통합 지점)까지 "
            "제공한다.",
            "근거: ETRI 2026-06-17 킥오프 회의록 - 상명대(서대희 교수·김원빈 연구교수·"
            "박건우·김도현) / ETRI(나갑주 팀장·김선형·정진욱) 참석.",
        ],
        body_size=13,
    )

    add_footer(slide, 2)


# ---------------------------------------------------------------------------
# 슬라이드 3: 1차년도 -> 2차년도 업그레이드 Before/After
# ---------------------------------------------------------------------------

BEFORE_AFTER_ROWS = [
    ("보안", "인증/인가 없음, RabbitMQ 평문 비밀번호 관행",
     "X.509 PKI + mTLS + JWT(RS256) + RBAC + JWS 서명 + 감사로그 전면 도입"),
    ("비밀 관리", "매니페스트에 비밀번호 하드코딩",
     "배포 스크립트가 openssl rand로 매 배포 시 Secret 생성, 저장소에 비밀값 없음"),
    ("통신 구조", "프로토타입 수준, 별도 AAA 계층 없음",
     "제어(KubeEdge 채널) / 데이터·인증(REST+mTLS+JWT) 하이브리드 분리"),
    ("네트워크", "고정 IP 하드코딩",
     "__CLOUD_IP__ 플레이스홀더 + 게이트웨이/사설IP 집선 구조(src/eam/gateway)"),
    ("성능 검증", "별도 벤치마크 없음",
     "bench/ N-스윕 실측 + M/M/c 외삽 모델로 1,000기 SLA 정량 검증"),
    ("시연", "K8s 환경 시연만",
     "로컬(클러스터 불필요) + K8s 실환경, 보안 적용 전/후 비교 데모 2경로"),
]


def build_before_after(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "1차년도 -> 2차년도 업그레이드 개요", "출처: README.md '1차년도 대비 변경 요약'")

    left, top = Inches(0.5), Inches(1.15)
    total_w = Inches(12.33)
    col_w = [Inches(1.7), Inches(5.15), Inches(5.48)]
    header_h = Inches(0.45)
    row_h = Inches(0.85)

    headers = ["구분", "1차년도 (prototype-y1/)", "2차년도 (본 repo)"]
    x = left
    for i, htext in enumerate(headers):
        add_rect(slide, x, top, col_w[i], header_h, fill=NAVY_DARK)
        tb = slide.shapes.add_textbox(x + Inches(0.08), top, col_w[i] - Inches(0.16), header_h)
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_top = 0
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = htext
        _set_run_font(run, size=13, bold=True, color=WHITE)
        x += col_w[i]

    y = top + header_h
    for r, (label, before, after) in enumerate(BEFORE_AFTER_ROWS):
        row_fill = WHITE if r % 2 == 0 else GRAY_ACCENT
        x = left
        cells = [label, before, after]
        sizes = [12, 11, 11]
        bolds = [True, False, False]
        for i, ctext in enumerate(cells):
            add_rect(slide, x, y, col_w[i], row_h, fill=row_fill, line=GRAY_LINE, line_width=Pt(0.5))
            tb = slide.shapes.add_textbox(x + Inches(0.08), y + Inches(0.03), col_w[i] - Inches(0.16), row_h - Inches(0.06))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = 0
            tf.margin_top = 0
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = ctext
            _set_run_font(run, size=sizes[i], bold=bolds[i], color=GRAY_BODY if i else NAVY_DARK)
            x += col_w[i]
        y += row_h

    add_footer(slide, 3)


# ---------------------------------------------------------------------------
# 슬라이드 4: 시스템 아키텍처 (4계층, 도형으로 직접)
# ---------------------------------------------------------------------------

def _add_arrow_connector(slide, x1, y1, x2, y2, *, color=NAVY_MID, width=Pt(2.25)):
    connector = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = STRAIGHT
    connector.line.color.rgb = color
    connector.line.width = width
    ln = connector.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return connector


def _add_box(slide, left, top, width, height, title, bullets, *, fill=NAVY_DARK, title_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = NAVY_LIGHT
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    _set_run_font(run, size=15, bold=True, color=title_color)
    for b in bullets:
        bp = tf.add_paragraph()
        bp.alignment = PP_ALIGN.LEFT
        brun = bp.add_run()
        brun.text = f"• {b}"
        _set_run_font(brun, size=10.5, bold=False, color=WHITE if fill != WHITE else GRAY_BODY)
    return shape


def _add_label(slide, left, top, width, height, text, *, size=10.5, color=NAVY_DARK, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size=size, bold=True, color=color)
    return box


def build_architecture(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "시스템 아키텍처 - 4계층 보안 흐름", "Device - Agent - Manager - Backend (출처: docs/05-kubeedge-integration.md)")

    box_top = Inches(2.15)
    box_h = Inches(2.1)

    device_box = _add_box(
        slide, Inches(0.5), box_top, Inches(2.35), box_h,
        "Device", ["온디바이스 센서", "(가상/실)"], fill=NAVY_MID,
    )
    agent_box = _add_box(
        slide, Inches(3.55), box_top, Inches(2.55), box_h,
        "Agent (엣지)", ["EdgeAgent", "enroll/get_token", "send_telemetry"], fill=NAVY_DARK,
    )
    manager_box = _add_box(
        slide, Inches(6.75), box_top, Inches(2.9), box_h,
        "Manager (클라우드 코어)", ["X.509 CA", "RBAC 인가", "SQLite Audit"], fill=NAVY_DARK,
    )
    backend_box = _add_box(
        slide, Inches(10.3), box_top, Inches(2.53), box_h,
        "Backend", ["대시보드/감사 소비자", "GET /audit, /devices"], fill=NAVY_MID,
    )

    arrow_y = box_top + Emu(int(box_h) // 2)
    _add_arrow_connector(slide, Inches(2.85), arrow_y, Inches(3.55), arrow_y)
    _add_arrow_connector(slide, Inches(6.10), arrow_y, Inches(6.75), arrow_y)
    _add_arrow_connector(slide, Inches(9.65), arrow_y, Inches(10.3), arrow_y)

    _add_label(slide, Inches(2.7), box_top - Inches(0.42), Inches(1.2), Inches(0.4), "X.509 enroll\n(CSR)")
    _add_label(slide, Inches(6.0), box_top - Inches(0.42), Inches(1.2), Inches(0.4), "mTLS+JWT / JWS")
    _add_label(slide, Inches(9.55), box_top - Inches(0.42), Inches(1.2), Inches(0.4), "audit 조회")

    # Gateway (사설IP 집선) - Agent 레이어 아래에 배치, Manager로 별도 업링크
    gw_top = Inches(4.85)
    gw_box = _add_box(
        slide, Inches(3.55), gw_top, Inches(2.55), Inches(1.3),
        "Gateway", ["사설IP 하위 Device 집선", "배치 업링크"], fill=NAVY_MID,
    )
    sub_box = add_rect(slide, Inches(0.5), gw_top, Inches(2.35), Inches(1.3), fill=GRAY_ACCENT, line=GRAY_LINE, line_width=Pt(0.75))
    sub_tf = sub_box.text_frame
    sub_tf.word_wrap = True
    sub_tf.margin_left = Inches(0.08)
    sub_p = sub_tf.paragraphs[0]
    sub_p.alignment = PP_ALIGN.CENTER
    sub_run = sub_p.add_run()
    sub_run.text = "SubDevice 다수\n(사설망, 비공인 IP)"
    _set_run_font(sub_run, size=11, bold=False, color=GRAY_BODY)

    _add_arrow_connector(slide, Inches(2.85), gw_top + Inches(0.65), Inches(3.55), gw_top + Inches(0.65))
    _add_arrow_connector(
        slide, Inches(4.8), gw_top, Inches(7.6), box_top + box_h,
        color=NAVY_LIGHT, width=Pt(1.5),
    )
    _add_label(slide, Inches(4.9), Inches(4.4), Inches(2.6), Inches(0.4), "{gateway_id, batch:[...]} - 단일 mTLS 업링크", size=10, color=NAVY_MID, align=PP_ALIGN.LEFT)

    add_footer(slide, 4)


# ---------------------------------------------------------------------------
# 슬라이드 5~11: 수행항목 1~7
# ---------------------------------------------------------------------------

ITEM_SLIDES = [
    dict(
        no=1,
        title="K8s(KubeEdge) 구동·성능 (1,000기 지연 검증)",
        requirement=[
            "K8s(KubeEdge) 환경에서 1,000기 규모 디바이스 운용 시 지연시간 검증",
        ],
        solution=[
            "N=10~200 실측(enroll/auth/telemetry) 후 M/M/c 대기행렬로 1,000기 외삽",
            "토큰 캐시·게이트웨이 집선·세션 재사용으로 인증 스톰 부하 자체를 완화",
        ],
        artifacts=[
            "코드: k8s/*.yaml, deploy/demo-setup-v2.sh, bench/run_bench.py, bench/model.py",
            "문서: docs/06-performance-plan.md, docs/perf/PERFORMANCE_REPORT.md",
            "테스트: tests/test_manifests.py, tests/test_bench_model.py",
        ],
        demo="deploy/demo-setup-v2.sh(K8s 실환경) / bench/run_bench.py && python bench/report.py(실측 재현)",
        status="구현완료 + 문서화 + 시연가능 (K8s 실환경은 Multipass 필요)",
    ),
    dict(
        no=2,
        title="미들웨어 통신 방식 결정",
        requirement=[
            "KubeEdge 채널 기반 vs 별도 프로토콜(HTTPS/AMQP) 통신 방식 결정",
        ],
        solution=[
            "하이브리드 채택 - 제어(파드 스케줄링)는 KubeEdge 채널,",
            "데이터·인증(등록/AAA/텔레메트리)은 별도 mTLS REST 채널로 분리",
        ],
        artifacts=[
            "코드: src/eam/manager/app.py(REST API), k8s/rabbitmq.yaml(선택적 대안, 미사용)",
            "문서: docs/01-comm-method-decision.md, docs/02-manager-coexistence.md",
            "테스트: tests/test_manager_api.py",
        ],
        demo="해당 없음 (설계 결정 문서)",
        status="구현완료(하이브리드 구조로 반영) + 문서화",
    ),
    dict(
        no=3,
        title="대규모(1,000기) 디바이스 인증 AAA",
        requirement=[
            "1,000기 규모 디바이스에 대한 인증·인가·감사(AAA) 체계 구축",
        ],
        solution=[
            "X.509 PKI + mTLS + Bearer JWT(RS256) + RBAC + JWS 페이로드 서명",
            "+ SQLite 감사로그(모든 인증/인가 이벤트 기록)",
        ],
        artifacts=[
            "코드: src/eam/common/{pki,jws,audit}.py, src/eam/manager/{app,ca,rbac,store,schemas}.py, "
            "src/eam/simulator/{fleet,vdevice}.py",
            "문서: docs/04-rnr-interface.md, docs/06-performance-plan.md",
            "테스트: test_pki/test_jws/test_audit/test_manager_api/test_rbac/test_simulator.py",
        ],
        demo="python -m eam.simulator.fleet --n 200(실측), demo/run_demo.py(AAA 흐름 시연)",
        status="구현완료 + 문서화 + 시연가능  |  담당: 박건우 · 홍승우",
        highlight=True,
    ),
    dict(
        no=4,
        title="KubeEdge 연동 구조 적용",
        requirement=[
            "KubeEdge CloudCore-EdgeCore 구조 위에 4계층 보안 흐름 적용",
        ],
        solution=[
            "제어 평면(CloudHub-EdgeHub)은 그대로 유지, 그 위에 독립적인",
            "REST+mTLS+JWT/JWS AAA 계층으로 Device-Agent-Manager-Backend 4계층 관통",
        ],
        artifacts=[
            "코드: k8s/{manager,agent-edge1,agent-edge2,gateway,namespace}.yaml, deploy/demo-setup-v2.sh",
            "문서: docs/05-kubeedge-integration.md (4계층 흐름도 mermaid 포함)",
            "테스트: tests/test_manifests.py",
        ],
        demo="deploy/demo-setup-v2.sh + demo/DEMO_SCENARIO.md §2",
        status="구현완료 + 문서화 + 시연가능 (K8s 실환경 Multipass 필요, 로컬은 즉시)",
    ),
    dict(
        no=5,
        title="네트워크 구성·주소 체계 (공인 IP)",
        requirement=[
            "공인 IP 직접 연결 vs 게이트웨이/사설 IP 네트워크 주소 체계 결정",
        ],
        solution=[
            "권고: 게이트웨이 + 사설 IP - 공격표면 축소, 등록단위 축소, 사설 IP 재사용,",
            "Manager 동시 인증 부하 완화(4축 비교 모두 우위)",
        ],
        artifacts=[
            "코드: src/eam/gateway/gateway.py, k8s/gateway.yaml",
            "문서: docs/03-network-addressing.md",
            "테스트: tests/test_gateway.py",
        ],
        demo="demo/DEMO_SCENARIO.md §2.2 (게이트웨이 집선 확인)",
        status="구현완료 + 문서화 + 시연가능",
    ),
    dict(
        no=6,
        title="보안 모듈-프레임워크 R&R",
        requirement=[
            "상명대 보안 코어 모듈과 ETRI 프레임워크 간 역할·책임(R&R) 및",
            "연동 인터페이스 정의",
        ],
        solution=[
            "상명대 = 보안 코어(EAM) 제공자, ETRI = KubeEdge 플랫폼·상위 앱 제공자",
            "Manager REST API(/api/v1/*) + 인증서/JWT 신뢰 루트만이 유일한 통합 지점",
        ],
        artifacts=[
            "코드: src/eam/manager/app.py(공개 계약), src/eam/manager/schemas.py",
            "문서: docs/04-rnr-interface.md (엔드포인트 표, 인증서 프로파일, JWT 클레임 규격)",
            "테스트: tests/test_manager_api.py (엔드포인트 계약 검증)",
        ],
        demo="해당 없음 (문서 산출물)",
        status="문서화 완료",
    ),
    dict(
        no=7,
        title="시연 (보안 적용 전·후 비교)",
        requirement=[
            "보안 적용 전 상태와 적용 후 상태를 대조하는 시연 시나리오 구성",
        ],
        solution=[
            "INSECURE_MODE 스위치로 동일 Manager를 두 번 기동해 전/후를 완전히 분리·재현",
            "감사(accounting)는 모드와 무관하게 항상 기록 -> 같은 로그에서 비교 가능",
        ],
        artifacts=[
            "코드: demo/run_demo.py, src/eam/manager/app.py(INSECURE_MODE)",
            "문서: demo/DEMO_SCENARIO.md",
            "테스트: tests/test_manifests.py (run_demo.py --fast exit 0 검증 포함)",
        ],
        demo="python demo/run_demo.py / --fast, K8s: kubectl set env ... INSECURE_MODE=true|false",
        status="구현완료 + 문서화 + 시연가능 (로컬 검증 완료: exit 0)",
    ),
]


def build_item_slide(prs, item, page_no):
    slide = add_slide(prs)
    badge = "★핵심" if item.get("highlight") else ""
    add_title_bar(slide, f"수행항목 {item['no']}. {item['title']} {badge}".rstrip())

    left = Inches(0.5)
    top0 = Inches(1.12)
    half_w = Inches(5.9)
    gap = Inches(0.13)

    add_section_block(
        slide, left, top0, half_w, Inches(1.85),
        "요구 (해야 할 부분)", item["requirement"], body_size=13,
    )
    add_section_block(
        slide, left + half_w + gap, top0, half_w, Inches(1.85),
        "해결방안 (대응)", item["solution"], body_size=13,
    )

    artifacts_top = top0 + Inches(1.85) + Inches(0.15)
    artifact_lines = list(item["artifacts"])
    artifact_lines.append(f"데모: {item['demo']}")
    add_section_block(
        slide, left, artifacts_top, Inches(12.33), Inches(2.35),
        "본 저장소 산출물 (실제 경로: 코드·문서·테스트)", artifact_lines, body_size=12,
    )

    status_top = artifacts_top + Inches(2.35) + Inches(0.15)
    status_box = add_rect(slide, left, status_top, Inches(12.33), Inches(0.55), fill=OK_NAVY)
    stf = status_box.text_frame
    stf.vertical_anchor = MSO_ANCHOR.MIDDLE
    stf.margin_left = Inches(0.12)
    sp = stf.paragraphs[0]
    srun = sp.add_run()
    srun.text = f"상태: {item['status']}"
    _set_run_font(srun, size=13, bold=True, color=WHITE)

    add_footer(slide, page_no)


# ---------------------------------------------------------------------------
# 슬라이드 12: 1,000기 성능 검증 결과
# ---------------------------------------------------------------------------

def build_performance(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "1,000기 성능 검증 결과", "실측(N=10~200) 기반 M/M/c 외삽 - 출처: docs/perf/PERFORMANCE_REPORT.md")

    img_top = Inches(1.15)
    img_w = Inches(5.7)
    if P95_CHART.exists():
        slide.shapes.add_picture(str(P95_CHART), Inches(0.5), img_top, width=img_w)
    if THROUGHPUT_CHART.exists():
        slide.shapes.add_picture(str(THROUGHPUT_CHART), Inches(6.45), img_top, width=img_w)

    caption_top = img_top + Inches(3.8) + Inches(0.05)
    add_textbox(
        slide, Inches(0.5), caption_top, img_w, Inches(0.3),
        [("auth-only p95 지연시간 vs N", 10, False, GRAY_BODY)], align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(6.45), caption_top, img_w, Inches(0.3),
        [("auth-only 처리량 vs N", 10, False, GRAY_BODY)], align=PP_ALIGN.CENTER,
    )

    concl_top = caption_top + Inches(0.35)
    add_section_block(
        slide, Inches(0.5), concl_top, Inches(12.33), Inches(1.75),
        "1,000기 외삽 결론 (M/M/c, 단일 워커 c=1)",
        [
            "실측 단일 워커 서비스율 μ = 23.76 auth/s (N=25 auth-only 버스트 최댓값)",
            "1,000기 60초 내 동시 재인증 스톰 -> 유틸라이제이션 ρ = 0.70 (안정)",
            "대기시간(Wq): 평균 99.0ms / p95 372.6ms / p99 599.6ms (목표 SLA p95<1.0s 충족)",
            "권장 레플리카 수: 1대 (단일 인스턴스로 목표 SLA 충족)",
            "참고: 실측 auth p95 @ N=200 = 534.1ms (측정 최대 N에서의 실측치)",
        ],
        body_size=12.5,
    )

    add_footer(slide, 12)


# ---------------------------------------------------------------------------
# 슬라이드 13: 시연 시나리오
# ---------------------------------------------------------------------------

def build_demo_slide(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "시연 시나리오 - 보안 적용 전·후 비교", "출처: demo/DEMO_SCENARIO.md (수행항목 4·7)")

    steps = [
        ("1단계 (BEFORE)", "Manager를 INSECURE_MODE=true로 기동 - 등록/인증 없이 임의 device_id로 텔레메트리 직접 POST",
         "HTTP 200 {status: accepted} - 누구나 데이터를 주입할 수 있는 상태"),
        ("2단계-a (AFTER, 거부)", "동일 Manager를 INSECURE_MODE=false로 재기동 후 똑같은 미인증 요청 재현",
         "HTTP 401 즉시 거부"),
        ("2단계-b (AFTER, 정상흐름)", "정식 디바이스: CSR 등록 -> 인증서 발급 -> Bearer JWT 발급 -> JWS 서명 텔레메트리",
         "정식 절차를 거친 디바이스만 통과"),
        ("3단계 (감사 확인)", "감사로그(audit log) tail 출력",
         "미인증 요청도 401로 즉시 감사되고, 정식 디바이스의 register/auth_success/telemetry_accept 이벤트가 모두 기록"),
    ]

    top = Inches(1.15)
    row_h = Inches(1.28)
    for i, (stage, action, point) in enumerate(steps):
        y = top + i * row_h
        add_rect(slide, Inches(0.5), y, Inches(2.2), row_h - Inches(0.08), fill=NAVY_MID)
        stage_tb = slide.shapes.add_textbox(Inches(0.58), y, Inches(2.04), row_h - Inches(0.08))
        stf = stage_tb.text_frame
        stf.word_wrap = True
        stf.vertical_anchor = MSO_ANCHOR.MIDDLE
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = stage
        _set_run_font(srun, size=12.5, bold=True, color=WHITE)

        add_rect(slide, Inches(2.78), y, Inches(9.9), row_h - Inches(0.08), fill=None, line=GRAY_LINE, line_width=Pt(0.75))
        body_tb = slide.shapes.add_textbox(Inches(2.9), y + Inches(0.05), Inches(9.7), row_h - Inches(0.18))
        btf = body_tb.text_frame
        btf.word_wrap = True
        bp = btf.paragraphs[0]
        brun = bp.add_run()
        brun.text = f"내용: {action}"
        _set_run_font(brun, size=12, bold=False, color=GRAY_BODY)
        bp2 = btf.add_paragraph()
        bp2.space_before = Pt(2)
        brun2 = bp2.add_run()
        brun2.text = f"촬영 포인트: {point}"
        _set_run_font(brun2, size=11.5, bold=False, color=NAVY_MID)

    footer_note_top = top + len(steps) * row_h + Inches(0.05)
    add_textbox(
        slide, Inches(0.5), footer_note_top, Inches(12.33), Inches(0.5),
        [("로컬: python demo/run_demo.py [--fast]  |  K8s: kubectl set env deployment/manager -n edge-auth INSECURE_MODE=true|false", 11, False, GRAY_BODY)],
    )

    add_footer(slide, 13)


# ---------------------------------------------------------------------------
# 슬라이드 14~17: 기능-화면 매핑 (실제 실행 화면 캡처)
# ---------------------------------------------------------------------------

def add_screen_picture(slide, img_path, left, top, max_w, max_h, caption=None):
    """이미지를 (max_w, max_h) 박스 안에 비율 유지로 배치하고 캡션을 단다."""
    from PIL import Image as PILImage  # Pillow는 matplotlib 의존성으로 이미 설치됨

    with PILImage.open(img_path) as im:
        w, h = im.size
    caption_h = Inches(0.38) if caption else 0
    scale = min(max_w / w, (max_h - caption_h) / h)
    pic_w, pic_h = int(w * scale), int(h * scale)
    x = left + int((max_w - pic_w) / 2)
    # 박스 안에서 (이미지+캡션)을 세로 중앙 정렬해 하단 여백 쏠림을 방지한다.
    y = top + int((max_h - caption_h - pic_h) / 2)
    slide.shapes.add_picture(str(img_path), x, y, width=pic_w, height=pic_h)
    if caption:
        add_textbox(
            slide, left, y + pic_h + Inches(0.06), max_w, Inches(0.32),
            [(caption, 10.5, False, GRAY_BODY)], align=PP_ALIGN.CENTER,
        )
    return pic_h


SCREEN_MAP_ROWS = [
    ("항목 ⑦·④  보안 전·후 비교 시연",
     "demo/run_demo.py 실행 콘솔 - INSECURE 모드 수용(200) vs 보안 모드 거부(401) + 감사로그 tail", "슬라이드 15"),
    ("항목 ③  디바이스 등록·인가 관리 (AAA)",
     "GET /api/v1/devices 관리 API 응답 - 등록 상태(approved/revoked)·인증서 시리얼·last_seen", "슬라이드 16"),
    ("항목 ③  Accounting 감사로그",
     "GET /api/v1/audit - 등록/인증 성공/위조 JWS 거부/폐기 후 인증 실패 실기록", "슬라이드 16"),
    ("항목 ⑥  연동 인터페이스 (R&R)",
     "Swagger UI(/docs) - Manager 전체 API 계약 (docs/04 인터페이스 명세와 1:1)", "슬라이드 17"),
    ("항목 ①·③  대규모(1,000기) 인증 검증",
     "가상 디바이스 일괄 등록·인증·전송 CLI(python -m eam.simulator.fleet) - 성능 차트는 슬라이드 12", "슬라이드 17"),
    ("항목 ①·④  실제 K8s 클러스터 구동",
     "Kubernetes Dashboard(edge-auth 네임스페이스)와 kubectl 실제 상태 - manager/agent/gateway 파드 Running", "슬라이드 18"),
]


def build_screen_map_overview(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "기능-화면 매핑 총괄",
                  "모든 캡처는 실제 플랫폼 화면 - Swagger UI·Kubernetes Dashboard·실제 CLI 출력 (ppt/capture_*.py 로 재생성)")

    top = Inches(1.2)
    row_h = Inches(0.88)
    for i, (feature, screen, ref) in enumerate(SCREEN_MAP_ROWS):
        y = top + i * row_h
        add_rect(slide, Inches(0.5), y, Inches(3.5), row_h - Inches(0.1), fill=NAVY_MID)
        ftb = slide.shapes.add_textbox(Inches(0.62), y, Inches(3.26), row_h - Inches(0.1))
        ftf = ftb.text_frame
        ftf.word_wrap = True
        ftf.vertical_anchor = MSO_ANCHOR.MIDDLE
        frun = ftf.paragraphs[0].add_run()
        frun.text = feature
        _set_run_font(frun, size=12.5, bold=True, color=WHITE)

        add_rect(slide, Inches(4.08), y, Inches(7.15), row_h - Inches(0.1),
                 fill=None, line=GRAY_LINE, line_width=Pt(0.75))
        stb = slide.shapes.add_textbox(Inches(4.2), y, Inches(6.9), row_h - Inches(0.1))
        stf = stb.text_frame
        stf.word_wrap = True
        stf.vertical_anchor = MSO_ANCHOR.MIDDLE
        srun = stf.paragraphs[0].add_run()
        srun.text = screen
        _set_run_font(srun, size=12, color=GRAY_BODY)

        add_rect(slide, Inches(11.31), y, Inches(1.52), row_h - Inches(0.1), fill=GRAY_ACCENT)
        rtb = slide.shapes.add_textbox(Inches(11.31), y, Inches(1.52), row_h - Inches(0.1))
        rtf = rtb.text_frame
        rtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        rp = rtf.paragraphs[0]
        rp.alignment = PP_ALIGN.CENTER
        rrun = rp.add_run()
        rrun.text = ref
        _set_run_font(rrun, size=12, bold=True, color=NAVY_DARK)

    note_top = top + len(SCREEN_MAP_ROWS) * row_h + Inches(0.1)
    add_textbox(
        slide, Inches(0.5), note_top, Inches(12.33), Inches(0.6),
        [("캡처 화면의 모든 데이터(디바이스·감사 이벤트·지연 수치)는 실제 시스템 실행 결과이며, "
          "docs/screens/ 에 원본 PNG로 저장되어 있다.", 11, False, GRAY_BODY)],
    )
    add_footer(slide, 14)


def build_screen_demo(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "화면 1. 보안 적용 전·후 비교 시연 (수행항목 ⑦·④)",
                  "demo/run_demo.py --fast 실제 실행 출력")
    add_screen_picture(
        slide, SCREEN_DEMO, Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.55),
        caption="BEFORE: 미인증 주입 HTTP 200 수용 -> AFTER: 동일 요청 401 거부, 정식 등록·인증 디바이스만 성공, 하단은 감사로그 tail",
    )
    add_footer(slide, 15)


def build_screen_admin(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "화면 2. 디바이스 관리·감사로그 조회 (수행항목 ③ AAA)",
                  "실행 중인 Manager의 Swagger UI에서 직접 호출한 실제 응답 - RBAC(operator/admin) 적용")
    half_w = Inches(6.05)
    add_screen_picture(
        slide, SCREEN_DEVICES, Inches(0.5), Inches(1.2), half_w, Inches(5.35),
        caption="GET /api/v1/devices - 등록 상태·인증서 시리얼·last_seen (revoked 상태 포함)",
    )
    add_screen_picture(
        slide, SCREEN_AUDIT, Inches(0.5) + half_w + Inches(0.23), Inches(1.2), half_w, Inches(5.35),
        caption="GET /api/v1/audit - 등록/인증 성공, 위조 JWS 거부, 폐기 후 인증 실패 실기록",
    )
    add_footer(slide, 16)


def build_screen_interface(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "화면 3. 연동 인터페이스·대규모 인증 시뮬레이터 (수행항목 ⑥·①·③)",
                  "Swagger UI(/docs) + 가상 디바이스 fleet CLI 실제 실행 출력")
    half_w = Inches(6.05)
    add_screen_picture(
        slide, SCREEN_SWAGGER, Inches(0.5), Inches(1.2), half_w, Inches(5.35),
        caption="Swagger UI - Manager 전체 API 계약 (docs/04 §2 인터페이스 명세와 1:1 대응)",
    )
    add_screen_picture(
        slide, SCREEN_FLEET, Inches(0.5) + half_w + Inches(0.23), Inches(1.2), half_w, Inches(5.35),
        caption="python -m eam.simulator.fleet - N기 일괄 등록·인증·전송, 단계별 p50/p95/p99 산출",
    )
    add_footer(slide, 17)


def build_screen_k8s(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "화면 4. 실제 Kubernetes 클러스터 구동 (수행항목 ①·④)",
                  "kind 기반 2노드 클러스터(etri-edge)에 본 저장소 k8s/*.yaml 을 그대로 배포한 실제 상태")
    half_w = Inches(6.05)
    add_screen_picture(
        slide, SCREEN_K8S_DASHBOARD, Inches(0.5), Inches(1.2), half_w, Inches(5.35),
        caption="Kubernetes Dashboard - edge-auth 네임스페이스, manager/agent×2/gateway 4개 Deployment Running",
    )
    add_screen_picture(
        slide, SCREEN_K8S_PODS, Inches(0.5) + half_w + Inches(0.23), Inches(1.2), half_w, Inches(5.35),
        caption="kubectl 실제 출력 - edge 라벨 노드에 agent/gateway 배치, manager NodePort 30443",
    )
    add_footer(slide, 18)


# ---------------------------------------------------------------------------
# 슬라이드 18: 향후 계획
# ---------------------------------------------------------------------------

def build_future(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "향후 계획", "3차년도 이관 항목 및 ETRI 협의 항목")

    add_section_block(
        slide, Inches(0.5), Inches(1.15), Inches(6.0), Inches(5.5),
        "3차년도 이관 / 확장 검토 항목",
        [
            "장애 인지·자동 복구 - 2차년도는 시연용 최소 수준(401 거부+감사)만 제공",
            "프레임워크(KubeEdge/ETRI 플랫폼) 전체 통합 - 3차년도 또는 타 기관 연계",
            "실네트워크/K8s 환경 재측정 (현재 μ=23.76 auth/s는 in-process ASGI 기준)",
            "uvicorn 다중 워커 + SQLite -> 다른 저장소 전환 (레플리카 확장 시 락 경합 대응)",
            "CA 계층 구조 - 다기관 연동 시 Intermediate CA 필요 여부 결정",
            "인증서/키 영속화 - emptyDir -> PVC 또는 KMS/HSM 연동",
            "RabbitMQ 등 비동기 채널 채택 여부, EdgeMesh 도입 검토",
            "auth/token 키 소지 증명(서명 nonce 챌린지) - 현재는 인증서 제시만 검증",
            "RBAC default-deny 전환 - 미등재 엔드포인트의 암묵적 허용 제거",
            "JWS freshness(재전송 방지) - 텔레메트리 서명에 타임스탬프/nonce 검증 추가",
        ],
        body_size=11,
    )

    add_section_block(
        slide, Inches(6.7), Inches(1.15), Inches(6.13), Inches(5.5),
        "ETRI 협의 필요 항목",
        [
            "mTLS 종단 위치 - 리버스 프록시/Ingress 실제 구성 확정",
            "다수 Manager 트래픽 라우팅 - site/group 샤딩 규칙의 실배포 토폴로지 매핑",
            "목표 SLA 재정의 - 현재 Wq p95<1.0s는 임의 설정값, 실제 요구사항 확정 필요",
            "1,000기 규모 실증 테스트베드 - 물리/가상 디바이스 자원 배정",
            "회의 참석: 상명대(서대희 교수·김원빈 연구교수·박건우·김도현) / "
            "ETRI(나갑주 팀장·김선형·정진욱)",
        ],
        body_size=12.5,
    )

    add_footer(slide, 19)


# ---------------------------------------------------------------------------
# 조립 + 검증
# ---------------------------------------------------------------------------

def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover(prs)                       # 1
    build_scope(prs)                       # 2
    build_before_after(prs)                # 3
    build_architecture(prs)                # 4
    for idx, item in enumerate(ITEM_SLIDES):  # 5~11
        build_item_slide(prs, item, page_no=5 + idx)
    build_performance(prs)                 # 12
    build_demo_slide(prs)                  # 13
    build_screen_map_overview(prs)         # 14
    build_screen_demo(prs)                 # 15
    build_screen_admin(prs)                # 16
    build_screen_interface(prs)            # 17
    build_screen_k8s(prs)                  # 18
    build_future(prs)                      # 19

    return prs


def verify(path: Path) -> None:
    # (c) zip 무결성
    with zipfile.ZipFile(path) as zf:
        bad = zf.testzip()
        assert bad is None, f"zip 무결성 오류: {bad}"

    # (b) 재로드 + 슬라이드 수/텍스트프레임 경계 검증
    prs = Presentation(str(path))
    slides = list(prs.slides)
    assert len(slides) == 19, f"슬라이드 수 불일치: {len(slides)} (기대: 19)"

    for si, slide in enumerate(slides, start=1):
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue  # 그룹 내부 등 위치 미지정 도형은 스킵
            left, top = shape.left, shape.top
            width = shape.width or 0
            height = shape.height or 0
            assert left >= -Emu(1000), f"슬라이드 {si} 도형이 왼쪽 경계를 벗어남: {shape.shape_id}"
            assert top >= -Emu(1000), f"슬라이드 {si} 도형이 위쪽 경계를 벗어남: {shape.shape_id}"
            assert left + width <= prs.slide_width + Emu(1000), (
                f"슬라이드 {si} 도형이 오른쪽 경계를 벗어남: {shape.shape_id} "
                f"(left+width={left + width}, slide_width={prs.slide_width})"
            )
            assert top + height <= prs.slide_height + Emu(1000), (
                f"슬라이드 {si} 도형이 아래쪽 경계를 벗어남: {shape.shape_id} "
                f"(top+height={top + height}, slide_height={prs.slide_height})"
            )

    # 수행항목 3(★핵심, 슬라이드 7)의 제목에 "★"가 정확히 1개만 등장하는지 검증
    # (title 문자열의 리터럴 ★핵심 + badge의 ★가 중복되는 회귀를 막기 위한 가드)
    item3_slide = slides[6]  # 0-indexed: 슬라이드 7
    title_shape = item3_slide.shapes[1]  # add_title_bar: [0]=배경 사각형, [1]=제목 텍스트박스
    title_text = title_shape.text_frame.paragraphs[0].text
    star_count = title_text.count("★")
    assert star_count == 1, (
        f"슬라이드 7 제목에 '★'가 {star_count}개 등장(기대: 1개) - title='{title_text}'"
    )
    assert "핵심" in title_text, f"슬라이드 7 제목에 '★핵심' 표기가 없음 - title='{title_text}'"

    print(f"검증 통과: {len(slides)}개 슬라이드, zip 무결성 OK, 모든 도형이 슬라이드 경계 내부에 위치, 슬라이드 7 제목 ★ 1개.")


def main():
    args = sys.argv[1:]
    do_build = "--verify" not in args
    do_verify = "--no-verify" not in args

    if do_build:
        PPT_DIR.mkdir(parents=True, exist_ok=True)
        prs = build_presentation()
        prs.save(str(OUTPUT_PATH))
        print(f"생성 완료: {OUTPUT_PATH}")

    if do_verify:
        verify(OUTPUT_PATH)


if __name__ == "__main__":
    main()
